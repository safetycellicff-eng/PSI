"""Builds a Plant Safety Inspection deck matching the original meeting format.

One slide per record: green title bar, a 2-column details table, site photos
below, and a category marker at the bottom-right corner. Slides are added to
template.pptx so they inherit the original master branding/background.
"""

import io

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

TITLE_GREEN = RGBColor(0x92, 0xD0, 0x50)
LABEL_RED = RGBColor(0xE8, 0x41, 0x2C)
MARKER_BLUE = RGBColor(0x5B, 0x9B, 0xD5)  # theme accent1 of the original deck
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_WIDTH = Inches(11.25)
SLIDE_HEIGHT = Inches(7.5)

# Photo area between the details table and the bottom edge of the slide.
PHOTO_LEFT = Inches(0.7)
PHOTO_TOP = Inches(2.85)
PHOTO_WIDTH = Inches(9.85)
PHOTO_HEIGHT = Inches(4.05)
PHOTO_GAP = Inches(0.15)

LABEL_WIDTH = Inches(1.09)
LABEL_HEIGHT = Inches(0.4)


def build_ppt(records, photos_by_id, template_path="template.pptx",
              heading="PLANT SAFETY INSPECTION"):
    """Return the .pptx file bytes for the given records.

    records: list of dicts keyed by storage.RECORD_HEADERS.
    photos_by_id: dict mapping record ID -> list of image bytes.
    """
    prs = Presentation(template_path)
    layout = _blank_layout(prs)
    for record in records:
        slide = prs.slides.add_slide(layout)
        _add_title(slide, heading, record.get("Status", "Pending"))
        _add_details_table(slide, record)
        _add_photos(slide, photos_by_id.get(record.get("ID"), []))
        _add_category_marker(slide, record.get("Category", "SV"))
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            return layout
    return prs.slide_layouts[-1]


def _add_title(slide, heading, status):
    box = slide.shapes.add_textbox(0, 0, SLIDE_WIDTH, Inches(0.47))
    box.fill.solid()
    box.fill.fore_color.rgb = TITLE_GREEN
    tf = box.text_frame
    tf.margin_top = Emu(36283)
    tf.margin_bottom = Emu(36283)
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER

    run = para.add_run()
    run.text = f"{heading} ("
    _style_run(run, "Arial Black", 19, color=BLACK)

    run = para.add_run()
    run.text = status
    _style_run(run, "Times New Roman", 19, bold=True, color=BLACK)

    run = para.add_run()
    run.text = ")"
    _style_run(run, "Arial Black", 19, color=BLACK)


def _add_details_table(slide, record):
    rows = [
        ("Description of Violation/ Hazard",
         record.get("Description of Violation/Hazard", ""), True),
        ("First appeared on", record.get("First Appeared On", ""), False),
        ("Action", record.get("Action By", ""), False),
        ("Remarks", record.get("Remarks", ""), False),
    ]
    shape = slide.shapes.add_table(
        len(rows), 2, Inches(0.98), Inches(0.76), Inches(9.76), Inches(1.86)
    )
    table = shape.table
    table.first_row = False
    table.horz_banding = False
    table.columns[0].width = Inches(2.03)
    table.columns[1].width = Inches(7.73)
    for height, row in zip((0.62, 0.28, 0.29, 0.63), table.rows):
        row.height = Inches(height)

    for row_no, (label, value, value_bold) in enumerate(rows):
        for col_no, (text, bold) in enumerate(((label, False), (value, value_bold))):
            cell = table.cell(row_no, col_no)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            _set_cell_borders(cell)
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.text = str(text)
            _style_run(run, "Calibri", 16, bold=bold, color=BLACK)


def _add_photos(slide, photo_set):
    """Place photos below the table.

    photo_set: {"before": [bytes], "after": [bytes]}. When both kinds exist,
    before photos go on the left half and after photos on the right, with red
    BEFORE/AFTER labels underneath — matching the original meeting format.
    """
    if isinstance(photo_set, list):  # tolerate a plain list of photos
        photo_set = {"before": photo_set, "after": []}
    before = list(photo_set.get("before", []))[:4]
    after = list(photo_set.get("after", []))[:4]
    if not before and not after:
        return

    if before and after:
        half_w = (PHOTO_WIDTH - PHOTO_GAP) / 2
        area_h = PHOTO_HEIGHT - LABEL_HEIGHT - Inches(0.1)
        right_x = PHOTO_LEFT + half_w + PHOTO_GAP
        _fill_area(slide, before[:2], PHOTO_LEFT, PHOTO_TOP, half_w, area_h)
        _fill_area(slide, after[:2], right_x, PHOTO_TOP, half_w, area_h)
        label_y = PHOTO_TOP + area_h + Inches(0.08)
        _add_label(slide, "BEFORE", PHOTO_LEFT + (half_w - LABEL_WIDTH) / 2, label_y)
        _add_label(slide, "AFTER", right_x + (half_w - LABEL_WIDTH) / 2, label_y)
    else:
        _fill_area(slide, before or after, PHOTO_LEFT, PHOTO_TOP,
                   PHOTO_WIDTH, PHOTO_HEIGHT)


def _fill_area(slide, photos, left, top, width, height):
    """Lay out up to 4 photos in the given area, preserving aspect ratios."""
    count = len(photos)
    if count == 0:
        return
    if count == 1:
        grid = [(0, 0, 1, 1)]
    elif count == 2:
        if width >= height:  # side by side in wide areas, stacked in columns
            grid = [(0, 0, 2, 1), (1, 0, 2, 1)]
        else:
            grid = [(0, 0, 1, 2), (0, 1, 1, 2)]
    else:
        grid = [(col, row, 2, 2) for row in range(2) for col in range(2)]

    for photo_bytes, (col, row, n_cols, n_rows) in zip(photos, grid):
        cell_w = (width - PHOTO_GAP * (n_cols - 1)) / n_cols
        cell_h = (height - PHOTO_GAP * (n_rows - 1)) / n_rows
        cell_x = left + col * (cell_w + PHOTO_GAP)
        cell_y = top + row * (cell_h + PHOTO_GAP)
        with Image.open(io.BytesIO(photo_bytes)) as img:
            aspect = img.width / img.height
        pic_w = cell_w
        pic_h = Emu(int(pic_w / aspect))
        if pic_h > cell_h:
            pic_h = cell_h
            pic_w = Emu(int(pic_h * aspect))
        slide.shapes.add_picture(
            io.BytesIO(photo_bytes),
            cell_x + Emu(int((cell_w - pic_w) / 2)),
            cell_y + Emu(int((cell_h - pic_h) / 2)),
            pic_w, pic_h,
        )


def _add_label(slide, text, left, top):
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, LABEL_WIDTH, LABEL_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LABEL_RED
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_top = 0
    tf.margin_bottom = 0
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    _style_run(run, "Calibri", 14, bold=True, color=WHITE)


def _add_category_marker(slide, category):
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(10.42), Inches(7.04), Inches(0.69), Inches(0.37)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = MARKER_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = category or "SV"
    _style_run(run, "Calibri Light", 15, bold=True, color=WHITE)


def _style_run(run, font_name, size_pt, bold=False, color=None):
    font = run.font
    font.name = font_name
    font.size = Pt(size_pt)
    font.bold = bold
    if color is not None:
        font.color.rgb = color


def _set_cell_borders(cell, width_emu=10157):
    """Apply thin solid black borders on all four sides of a table cell."""
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for existing in tcPr.findall(qn(tag)):
            tcPr.remove(existing)
    # Border elements must appear before any fill element in tcPr.
    fill = tcPr.find(qn("a:solidFill"))
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        ln = tcPr.makeelement(qn(tag), {
            "w": str(width_emu), "cap": "flat", "cmpd": "sng", "algn": "ctr",
        })
        solid = ln.makeelement(qn("a:solidFill"), {})
        clr = ln.makeelement(qn("a:srgbClr"), {"val": "000000"})
        solid.append(clr)
        ln.append(solid)
        dash = ln.makeelement(qn("a:prstDash"), {"val": "solid"})
        ln.append(dash)
        if fill is not None:
            fill.addprevious(ln)
        else:
            tcPr.append(ln)
